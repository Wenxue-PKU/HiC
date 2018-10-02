import datetime
import argparse
from pptx import Presentation
from pptx.util import Cm
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

parser = argparse.ArgumentParser()
parser.add_argument("-d", "--directory", help="fastqc directory")
parser.add_argument("-n", "--name", help="sample name(s) separated by commma,")
parser.add_argument("-o", "--out", help="output powerpoint file name")
args = parser.parse_args()

DIR = args.directory
FILE_OUT = args.out

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
title_and_content_layout = prs.slide_layouts[1]
title_only = prs.slide_layouts[5]
blank_slide_layout = prs.slide_layouts[6]


for NAME in args.name.split(','):
    FILE_ZIP = DIR + '/' + NAME + '_fastqc.zip'
    os.system('unzip -q ' + FILE_ZIP + ' -d ' + DIR + '/' + NAME)
    DIR_image = DIR + '/' + NAME + '/' + NAME + '_fastqc/Images/'
    img_base_quality = DIR_image + 'per_base_quality.png'
    img_content = DIR_image + 'per_base_sequence_content.png'
    img_duplication = DIR_image + 'duplication_levels.png'
    img_gc = DIR_image + 'per_sequence_gc_content.png'
    img_seq_quality = DIR_image + 'per_sequence_quality.png'
    
    
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.add_textbox(Cm(0.44), Cm(0.38), Cm(12.34), Cm(1.11))
    p = title.text_frame.paragraphs[0].add_run()
    p.text = 'Sequence quality of ' + NAME
    p.font.name = 'Segoe UI'
    p.font.size = Pt(24)
    p.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    pic = slide.shapes.add_picture(img_base_quality, Cm(0.87), Cm(2.12), height=Cm(8.18))
    pic = slide.shapes.add_picture(img_content, Cm(0.87), Cm(10.3), height=Cm(7.97))
    pic = slide.shapes.add_picture(img_duplication, Cm(16.83), Cm(0.77), height=Cm(5.63))
    pic = slide.shapes.add_picture(img_gc, Cm(16.83), Cm(6.81), height=Cm(5.63))
    pic = slide.shapes.add_picture(img_seq_quality, Cm(16.83), Cm(12.85), height=Cm(5.63))

    def add_text(word, left, top, height, width):
   	comment = slide.shapes.add_textbox(left, top, width, height)
   	p = comment.text_frame.paragraphs[0].add_run()
   	p.text = word
   	p.font.name = 'Segoe UI'
   	p.font.size = Pt(14)
    
    add_text('Quality score across all bases', Cm(1.19), Cm(8.9), Cm(1.03), Cm(9.3))
    add_text('Sequence content across all bases', Cm(7.38), Cm(10.87), Cm(1.03), Cm(9.47))
    add_text('Duplication level', Cm(18.69), Cm(0.09), Cm(1.03), Cm(5.04))
    add_text('Mean GC content', Cm(18.59), Cm(6.33), Cm(1.03), Cm(5.04))
    add_text('Mean sequence quality', Cm(17.97), Cm(12.38), Cm(1.03), Cm(6.65))
    
    os.system('rm -r ' + DIR + '/' + NAME)

prs.save(FILE_OUT)


